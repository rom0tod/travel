"""Генератор презентации к защите проекта TripPlanner.

Запуск:
    python -m tools.build_presentation

Создаёт ``docs/TripPlanner_presentation.pptx`` с фирменной палитрой
сайта (синий → фиолетовый градиент, оранжевый акцент).
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# --- Палитра -----------------------------------------------------------------

BLUE = RGBColor(0x25, 0x63, 0xEB)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
NAVY = RGBColor(0x0F, 0x17, 0x2A)
NAVY_SOFT = RGBColor(0x1E, 0x29, 0x3B)
BG = RGBColor(0xF8, 0xFA, 0xFC)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
TEXT = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x64, 0x74, 0x8B)
ACCENT = RGBColor(0xF9, 0x73, 0x16)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HEADER_FONT = "Calibri"
BODY_FONT = "Calibri"


# --- Утилиты -----------------------------------------------------------------

def add_filled_rect(slide, left, top, width, height,
                    fill_color, *, line_color=None):
    """Прямоугольник без обводки (по умолчанию)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_rounded_card(slide, left, top, width, height,
                     fill_color=CARD, border_color=BORDER):
    """Карточка со скруглёнными углами и тонкой границей."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text,
             *, font=BODY_FONT, size=14, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    """Текстовое поле с настроенным внутренним отступом."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_rich_text(slide, left, top, width, height, paragraphs,
                  *, anchor=MSO_ANCHOR.TOP, line_spacing=1.2,
                  align=PP_ALIGN.LEFT):
    """Текст с несколькими абзацами, каждый со своими опциями.

    paragraphs: list[dict] с ключами: text, size, bold, color, font,
    space_after (Pt), bullet (bool).
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    for index, spec in enumerate(paragraphs):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", align)
        p.line_spacing = spec.get("line_spacing", line_spacing)
        if "space_after" in spec:
            p.space_after = Pt(spec["space_after"])
        run = p.add_run()
        run.text = ("• " if spec.get("bullet") else "") + spec["text"]
        run.font.name = spec.get("font", BODY_FONT)
        run.font.size = Pt(spec.get("size", 14))
        run.font.bold = spec.get("bold", False)
        run.font.color.rgb = spec.get("color", TEXT)
    return box


def set_slide_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, text, dark=False):
    color = MUTED if not dark else RGBColor(0xCB, 0xD5, 0xE1)
    add_text(
        slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
        text, size=10, color=color, align=PP_ALIGN.LEFT,
    )


def add_page_number(slide, number, total, dark=False):
    color = MUTED if not dark else RGBColor(0xCB, 0xD5, 0xE1)
    add_text(
        slide, Inches(12.0), Inches(7.05), Inches(1.0), Inches(0.3),
        f"{number} / {total}", size=10, color=color,
        align=PP_ALIGN.RIGHT,
    )


# --- Слайды ------------------------------------------------------------------

TOTAL_SLIDES = 12


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide, NAVY)

    # Боковая акцентная полоса
    add_filled_rect(slide, Inches(0), Inches(0), Inches(0.3),
                    Inches(7.5), ACCENT)

    # Лейбл
    add_text(
        slide, Inches(1.0), Inches(1.5), Inches(11), Inches(0.4),
        "ЯНДЕКС ЛИЦЕЙ · WEBSERVER + REST API",
        font=BODY_FONT, size=14, bold=True,
        color=RGBColor(0xF9, 0x73, 0x16),
    )

    # Главный заголовок
    add_text(
        slide, Inches(1.0), Inches(2.0), Inches(11), Inches(1.5),
        "TripPlanner",
        font=HEADER_FONT, size=72, bold=True, color=WHITE,
    )

    # Подзаголовок
    add_text(
        slide, Inches(1.0), Inches(3.4), Inches(11), Inches(1.2),
        "Социальный планировщик путешествий\n"
        "с интерактивной картой и REST API",
        font=HEADER_FONT, size=26,
        color=RGBColor(0xCB, 0xD5, 0xE1),
        line_spacing=1.25,
    )

    # Автор
    add_text(
        slide, Inches(1.0), Inches(5.6), Inches(11), Inches(0.5),
        "Автор: Роман Тодоров",
        font=BODY_FONT, size=18, bold=True, color=WHITE,
    )
    add_text(
        slide, Inches(1.0), Inches(6.05), Inches(11), Inches(0.4),
        "Итоговый проект курса «Основы промышленного программирования» · 2026",
        font=BODY_FONT, size=14, color=RGBColor(0x94, 0xA3, 0xB8),
    )


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG)

    _add_section_header(slide, "01", "Идея и проблема")

    # Левая колонка
    add_text(
        slide, Inches(0.7), Inches(1.7), Inches(7.5), Inches(0.6),
        "Планировать поездку — больно",
        size=30, bold=True, color=NAVY,
    )
    add_rich_text(
        slide, Inches(0.7), Inches(2.5), Inches(7.5), Inches(4),
        [
            {"text": "Маршруты разбросаны по заметкам и скриншотам.",
             "size": 16, "color": TEXT, "space_after": 12},
            {"text": "Бюджет считается «на салфетке».",
             "size": 16, "color": TEXT, "space_after": 12},
            {"text": "Поделиться удобной картой с друзьями — отдельный квест.",
             "size": 16, "color": TEXT, "space_after": 12},
            {"text": "Готовые сервисы либо платные, либо привязаны к "
                     "экосистеме (Google Maps, Booking).",
             "size": 16, "color": TEXT},
        ],
    )

    # Правая колонка: цитата-плакат
    add_filled_rect(slide, Inches(8.6), Inches(1.7), Inches(4.0),
                    Inches(4.5), BLUE)
    add_text(
        slide, Inches(8.8), Inches(2.0), Inches(3.6), Inches(0.6),
        "ЦЕЛЬ", size=12, bold=True,
        color=RGBColor(0xCB, 0xD5, 0xE1),
    )
    add_text(
        slide, Inches(8.8), Inches(2.6), Inches(3.6), Inches(3.5),
        "Один сервис,\n"
        "где маршрут,\n"
        "бюджет и социалка\n"
        "живут вместе.",
        size=22, bold=True, color=WHITE, line_spacing=1.3,
    )

    add_footer(slide, "TripPlanner · защита проекта")
    add_page_number(slide, 2, TOTAL_SLIDES)


def slide_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG)
    _add_section_header(slide, "02", "Что умеет приложение")

    # Сетка 2x3 карточек
    cards = [
        ("Аккаунты",
         "Регистрация, хэширование паролей, профили с аватарами."),
        ("Поездки",
         "Создание, редактирование, теги, обложки, флаг публичности."),
        ("Маршрут",
         "Интерактивная карта Leaflet, точки по дням, 6 категорий."),
        ("Геокодер",
         "Авто-поиск координат адреса через Nominatim (сторонний API)."),
        ("Бюджет",
         "Расходы по 6 категориям, сводка, индикатор перерасхода."),
        ("Социалка",
         "Лайки (AJAX), комментарии, explore с поиском и пагинацией."),
    ]

    card_w = Inches(4.0)
    card_h = Inches(1.55)
    start_x = Inches(0.5)
    start_y = Inches(1.7)
    gap_x = Inches(0.15)
    gap_y = Inches(0.2)

    for i, (title, body) in enumerate(cards):
        col = i % 3
        row = i // 3
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        add_rounded_card(slide, x, y, card_w, card_h)
        # Цветная цифра
        add_text(
            slide, x + Inches(0.2), y + Inches(0.15),
            Inches(0.6), Inches(0.5),
            f"{i + 1:02d}", size=18, bold=True, color=BLUE,
        )
        add_text(
            slide, x + Inches(0.85), y + Inches(0.15),
            card_w - Inches(1.0), Inches(0.45),
            title, size=18, bold=True, color=NAVY,
        )
        add_text(
            slide, x + Inches(0.2), y + Inches(0.7),
            card_w - Inches(0.4), card_h - Inches(0.85),
            body, size=12, color=MUTED, line_spacing=1.3,
        )

    # Итоговая полоса
    add_filled_rect(slide, Inches(0.5), Inches(5.45), Inches(12.3),
                    Inches(1.0), NAVY)
    add_text(
        slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.8),
        "Полноценный продукт, а не учебная заготовка: каждая функция "
        "доступна и через UI, и через REST API.",
        size=15, bold=True, color=WHITE, line_spacing=1.3,
    )

    add_footer(slide, "TripPlanner · защита проекта")
    add_page_number(slide, 3, TOTAL_SLIDES)


def slide_journey(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG)
    _add_section_header(slide, "03", "Сценарий пользователя")

    steps = [
        ("Регистрация",
         "Логин, email,\nхэш пароля"),
        ("Создание\nпоездки",
         "Название, даты,\nбюджет, теги"),
        ("Точки\nна карте",
         "Geocoder или\nклик по карте"),
        ("Учёт\nбюджета",
         "Категории трат,\nсводка"),
        ("Публикация",
         "Лайки,\nкомментарии"),
    ]

    step_w = Inches(2.2)
    step_h = Inches(3.0)
    gap = Inches(0.25)
    total_w = step_w * len(steps) + gap * (len(steps) - 1)
    start_x = (Inches(13.33) - total_w) / 2
    start_y = Inches(2.0)

    for index, (title, body) in enumerate(steps):
        x = start_x + index * (step_w + gap)

        # Круг с номером сверху
        circle_size = Inches(0.7)
        circle_x = x + (step_w - circle_size) / 2
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, circle_x, start_y, circle_size, circle_size,
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT if index == 2 else BLUE
        circle.line.fill.background()
        circle.shadow.inherit = False
        add_text(
            slide, circle_x, start_y, circle_size, circle_size,
            str(index + 1), size=20, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Карточка
        card_y = start_y + circle_size - Inches(0.05)
        add_rounded_card(slide, x, card_y, step_w, step_h - circle_size)
        add_text(
            slide, x + Inches(0.1), card_y + Inches(0.35),
            step_w - Inches(0.2), Inches(0.9),
            title, size=16, bold=True, color=NAVY,
            align=PP_ALIGN.CENTER, line_spacing=1.15,
        )
        add_text(
            slide, x + Inches(0.1), card_y + Inches(1.25),
            step_w - Inches(0.2), Inches(1.1),
            body, size=11, color=MUTED,
            align=PP_ALIGN.CENTER, line_spacing=1.3,
        )

    add_text(
        slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.5),
        "От пустого аккаунта до публикации маршрута — 5–7 минут.",
        size=14, color=MUTED, align=PP_ALIGN.CENTER,
    )

    add_footer(slide, "TripPlanner · защита проекта")
    add_page_number(slide, 4, TOTAL_SLIDES)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG)
    _add_section_header(slide, "04", "Архитектура")

    add_text(
        slide, Inches(0.7), Inches(1.7), Inches(8), Inches(0.5),
        "Слойная архитектура · фабрика приложения",
        size=18, bold=True, color=NAVY,
    )

    # Слои
    layers = [
        ("Представление",
         "templates/  ·  static/  ·  Jinja2  ·  Bootstrap 5  ·  Leaflet",
         BLUE),
        ("Контроллеры",
         "main.py (Flask routes)  +  api/* (Flask-RESTful)",
         RGBColor(0x06, 0x5A, 0x82)),
        ("Бизнес-логика",
         "forms/  (WTForms)   tools/  (helpers, geocoder)",
         PURPLE),
        ("Данные",
         "data/  ·  SQLAlchemy 2.0  ·  PostgreSQL / SQLite",
         RGBColor(0x4C, 0x1D, 0x95)),
    ]

    layer_x = Inches(0.7)
    layer_y = Inches(2.3)
    layer_w = Inches(7.2)
    layer_h = Inches(0.95)
    gap = Inches(0.15)

    for index, (name, body, color) in enumerate(layers):
        y = layer_y + index * (layer_h + gap)
        add_filled_rect(slide, layer_x, y, layer_w, layer_h, color)
        add_text(
            slide, layer_x + Inches(0.3), y + Inches(0.12),
            Inches(2.5), Inches(0.4),
            name, size=15, bold=True, color=WHITE,
        )
        add_text(
            slide, layer_x + Inches(0.3), y + Inches(0.5),
            layer_w - Inches(0.6), Inches(0.4),
            body, size=12, color=RGBColor(0xE2, 0xE8, 0xF0),
        )

    # Правая колонка: ключевые принципы
    side_x = Inches(8.3)
    add_rounded_card(slide, side_x, Inches(2.3), Inches(4.5),
                     Inches(4.55))
    add_text(
        slide, side_x + Inches(0.3), Inches(2.45),
        Inches(4.0), Inches(0.4),
        "ПРИНЦИПЫ", size=11, bold=True, color=ACCENT,
    )
    add_rich_text(
        slide, side_x + Inches(0.3), Inches(2.85),
        Inches(4.0), Inches(4.0),
        [
            {"text": "Фабрика create_app()",
             "size": 14, "bold": True, "color": NAVY, "space_after": 4},
            {"text": "Изоляция dev / prod, удобно для тестов.",
             "size": 12, "color": MUTED, "space_after": 12},
            {"text": "Слой данных независим от Flask",
             "size": 14, "bold": True, "color": NAVY, "space_after": 4},
            {"text": "Модели работают в сидере и CLI без HTTP.",
             "size": 12, "color": MUTED, "space_after": 12},
            {"text": "Единая конфигурация",
             "size": 14, "bold": True, "color": NAVY, "space_after": 4},
            {"text": "Все константы — в config.py + ENV.",
             "size": 12, "color": MUTED, "space_after": 12},
            {"text": "Переиспользуемые шаблоны",
             "size": 14, "bold": True, "color": NAVY, "space_after": 4},
            {"text": "_card.html включается в 4 экранах.",
             "size": 12, "color": MUTED},
        ],
    )

    add_footer(slide, "TripPlanner · защита проекта")
    add_page_number(slide, 5, TOTAL_SLIDES)


def slide_data_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG)
    _add_section_header(slide, "05", "Модель данных")

    add_text(
        slide, Inches(0.7), Inches(1.7), Inches(11), Inches(0.5),
        "7 ORM-моделей · one-to-many и many-to-many · UNIQUE на лайках",
        size=16, color=MUTED,
    )

    # Центральная карточка User → Trip
    user_box = (Inches(0.8), Inches(2.5), Inches(2.2), Inches(0.9))
    trip_box = (Inches(4.5), Inches(2.5), Inches(2.2), Inches(0.9))

    def entity(x, y, w, h, name, fields, color=BLUE):
        add_rounded_card(slide, x, y, w, h)
        add_filled_rect(slide, x, y, w, Inches(0.35), color)
        add_text(slide, x + Inches(0.15), y + Inches(0.03),
                 w, Inches(0.3),
                 name, size=13, bold=True, color=WHITE)
        add_text(slide, x + Inches(0.15), y + Inches(0.4),
                 w - Inches(0.2), h - Inches(0.4),
                 fields, size=10, color=MUTED, line_spacing=1.2)

    entity(*user_box, "User",
           "username · email\npassword_hash\nabout · avatar")
    entity(*trip_box, "Trip",
           "title · destination\nstart/end · budget\nis_public · tags",
           color=PURPLE)
    entity(Inches(8.2), Inches(2.5), Inches(2.2), Inches(0.9),
           "Place",
           "name · kind · day\nlatitude · longitude\norder_index",
           color=BLUE)
    entity(Inches(10.7), Inches(2.5), Inches(2.2), Inches(0.9),
           "Expense",
           "category · amount\ncurrency · spent_at",
           color=BLUE)

    entity(Inches(4.5), Inches(4.4), Inches(2.2), Inches(0.9),
           "Comment",
           "content · created_at\nFK Trip + FK User",
           color=RGBColor(0x06, 0x5A, 0x82))
    entity(Inches(7.0), Inches(4.4), Inches(2.2), Inches(0.9),
           "Like",
           "FK Trip + FK User\nUNIQUE constraint",
           color=RGBColor(0x06, 0x5A, 0x82))
    entity(Inches(9.5), Inches(4.4), Inches(2.2), Inches(0.9),
           "Tag",
           "name · slug\nM2M ↔ Trip",
           color=RGBColor(0x06, 0x5A, 0x82))

    # Подпись внизу
    add_filled_rect(slide, Inches(0.7), Inches(5.7), Inches(12.0),
                    Inches(0.9), NAVY)
    add_text(
        slide, Inches(0.9), Inches(5.82), Inches(11.7), Inches(0.7),
        "Все модели наследуются от SerializerMixin → единый формат "
        "to_dict() для REST API. Каскадное удаление сохраняет "
        "целостность данных.",
        size=13, color=WHITE, line_spacing=1.3,
    )

    add_footer(slide, "TripPlanner · защита проекта")
    add_page_number(slide, 6, TOTAL_SLIDES)


def slide_tech_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG)
    _add_section_header(slide, "06", "Технологии")

    columns = [
        ("Backend", BLUE, [
            "Python 3.12",
            "Flask 2.3",
            "SQLAlchemy 2.0",
            "Flask-Login",
            "Flask-WTF",
            "Flask-RESTful",
            "Werkzeug security",
        ]),
        ("Frontend", PURPLE, [
            "Jinja2",
            "Bootstrap 5.3",
            "Bootstrap Icons",
            "Leaflet 1.9",
            "OpenStreetMap",
            "AJAX (fetch API)",
            "Адаптивная вёрстка",
        ]),
        ("Данные · DevOps", ACCENT, [
            "SQLite (dev)",
            "PostgreSQL (prod)",
            "sqlalchemy-serializer",
            "Nominatim API",
            "gunicorn (WSGI)",
            "Render.com",
            "GitHub + git",
        ]),
    ]

    col_w = Inches(4.0)
    col_h = Inches(4.8)
    gap = Inches(0.2)
    start_x = Inches(0.5)
    start_y = Inches(1.7)

    for col_index, (title, color, items) in enumerate(columns):
        x = start_x + col_index * (col_w + gap)
        add_rounded_card(slide, x, start_y, col_w, col_h)
        add_filled_rect(slide, x, start_y, col_w, Inches(0.6), color)
        add_text(
            slide, x + Inches(0.3), start_y + Inches(0.12),
            col_w - Inches(0.4), Inches(0.4),
            title, size=18, bold=True, color=WHITE,
        )
        for row_index, item in enumerate(items):
            y = start_y + Inches(0.85) + row_index * Inches(0.5)
            # Маленький круг-маркер
            bullet_size = Inches(0.18)
            bullet = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x + Inches(0.3), y + Inches(0.1),
                bullet_size, bullet_size,
            )
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = color
            bullet.line.fill.background()
            bullet.shadow.inherit = False

            add_text(
                slide, x + Inches(0.6), y, col_w - Inches(0.8),
                Inches(0.4),
                item, size=14, color=TEXT,
            )

    # Маркер «все технологии модуля закрыты»
    add_filled_rect(slide, Inches(0.5), Inches(6.6), Inches(12.3),
                    Inches(0.55), NAVY)
    add_text(
        slide, Inches(0.7), Inches(6.66), Inches(12.0), Inches(0.5),
        "✓ Все технологии модуля закрыты: requirements.txt · Bootstrap "
        "· ORM · auth · файлы · REST · сторонний API · БД · хостинг",
        size=13, bold=True, color=WHITE,
    )

    add_footer(slide, "TripPlanner · защита проекта")
    add_page_number(slide, 7, TOTAL_SLIDES)


def slide_api(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG)
    _add_section_header(slide, "07", "REST API")

    # Левая часть — список endpoint'ов
    add_text(
        slide, Inches(0.7), Inches(1.7), Inches(7), Inches(0.5),
        "14 эндпоинтов, единый формат ответа",
        size=16, bold=True, color=NAVY,
    )

    endpoints = [
        ("GET",    "/api/trips", "список публичных поездок"),
        ("POST",   "/api/trips", "создать поездку"),
        ("GET",    "/api/trips/<id>", "детали + точки и расходы"),
        ("PUT",    "/api/trips/<id>", "обновить (владелец)"),
        ("DELETE", "/api/trips/<id>", "удалить"),
        ("GET",    "/api/trips/<id>/places", "точки маршрута"),
        ("POST",   "/api/trips/<id>/places", "добавить точку"),
        ("GET",    "/api/trips/<id>/expenses", "расходы и сводка"),
        ("GET",    "/api/users/<id>", "профиль + публичные поездки"),
        ("GET",    "/api/geocode?address=", "Nominatim прокси"),
    ]

    method_colors = {
        "GET": RGBColor(0x16, 0xA3, 0x4A),
        "POST": BLUE,
        "PUT": ACCENT,
        "DELETE": RGBColor(0xDC, 0x26, 0x26),
    }

    list_x = Inches(0.7)
    list_y = Inches(2.3)
    row_h = Inches(0.43)
    list_w = Inches(7.3)

    for index, (method, path, comment) in enumerate(endpoints):
        y = list_y + index * row_h
        # Цветной бейдж метода
        add_filled_rect(
            slide, list_x, y + Inches(0.04),
            Inches(0.85), Inches(0.32),
            method_colors[method],
        )
        add_text(
            slide, list_x, y + Inches(0.04),
            Inches(0.85), Inches(0.32),
            method, size=10, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        # Путь
        add_text(
            slide, list_x + Inches(1.0), y + Inches(0.02),
            Inches(3.6), Inches(0.4),
            path, font="Consolas", size=12, color=NAVY,
        )
        # Комментарий
        add_text(
            slide, list_x + Inches(4.5), y + Inches(0.02),
            list_w - Inches(4.5), Inches(0.4),
            comment, size=11, color=MUTED,
        )

    # Правая часть — пример JSON-ответа
    code_x = Inches(8.4)
    add_filled_rect(slide, code_x, Inches(1.7), Inches(4.5),
                    Inches(5.0), NAVY)
    add_text(
        slide, code_x + Inches(0.25), Inches(1.85),
        Inches(4.0), Inches(0.4),
        "GET /api/trips", font="Consolas", size=14, bold=True,
        color=RGBColor(0xCB, 0xD5, 0xE1),
    )
    add_text(
        slide, code_x + Inches(0.25), Inches(2.3),
        Inches(4.0), Inches(4.3),
        '{\n'
        '  "success": true,\n'
        '  "total": 3,\n'
        '  "limit": 20,\n'
        '  "trips": [\n'
        '    {\n'
        '      "id": 2,\n'
        '      "title": "Гастротур\n'
        '         по Тбилиси",\n'
        '      "duration_days": 6,\n'
        '      "likes_count": 1,\n'
        '      "places_count": 4\n'
        '    }\n'
        '  ]\n'
        '}',
        font="Consolas", size=11, color=WHITE, line_spacing=1.25,
    )

    add_footer(slide, "TripPlanner · защита проекта")
    add_page_number(slide, 8, TOTAL_SLIDES)


def slide_highlights(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG)
    _add_section_header(slide, "08", "Технические находки")

    items = [
        ("Карта Leaflet",
         "Цветные маркеры + линия маршрута по дням, "
         "выбор точек кликом."),
        ("Геокодер Nominatim",
         "Серверный proxy с собственным User-Agent и "
         "обработкой ошибок сети."),
        ("AJAX-лайки",
         "Кнопка обновляет счётчик через fetch без "
         "перезагрузки страницы."),
        ("Авто-сидинг БД",
         "При старте на пустой БД сам создаёт демо-данные "
         "(нужно для бесплатного Render без Shell)."),
        ("Multi-уровневые права",
         "Trip.is_visible_to() / is_editable_by() — общий "
         "контракт для view и API."),
        ("One-click деплой",
         "render.yaml поднимает веб-сервис + PostgreSQL "
         "одной кнопкой."),
    ]

    card_w = Inches(4.0)
    card_h = Inches(2.3)
    gap_x = Inches(0.15)
    gap_y = Inches(0.2)
    start_x = Inches(0.5)
    start_y = Inches(1.7)

    for i, (title, body) in enumerate(items):
        col = i % 3
        row = i // 3
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        add_rounded_card(slide, x, y, card_w, card_h)
        # Толстая полоса слева как акцент
        add_filled_rect(slide, x, y, Inches(0.12), card_h,
                        ACCENT if i % 2 else BLUE)

        add_text(
            slide, x + Inches(0.4), y + Inches(0.2),
            card_w - Inches(0.5), Inches(0.5),
            title, size=18, bold=True, color=NAVY,
        )
        add_text(
            slide, x + Inches(0.4), y + Inches(0.85),
            card_w - Inches(0.5), card_h - Inches(1.0),
            body, size=12, color=TEXT, line_spacing=1.4,
        )

    add_footer(slide, "TripPlanner · защита проекта")
    add_page_number(slide, 9, TOTAL_SLIDES)


def slide_numbers(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, NAVY)
    _add_section_header(slide, "09", "В цифрах", dark=True)

    stats = [
        ("4 410", "строк кода", "Python, HTML, CSS, JS"),
        ("7", "ORM-моделей", "со связями one-to-many и M2M"),
        ("14", "REST API", "endpoints в едином формате"),
        ("30", "веб-маршрутов", "Flask routes"),
    ]

    card_w = Inches(2.9)
    card_h = Inches(3.0)
    gap = Inches(0.25)
    total_w = card_w * 4 + gap * 3
    start_x = (Inches(13.33) - total_w) / 2
    y = Inches(2.6)

    for i, (number, label, hint) in enumerate(stats):
        x = start_x + i * (card_w + gap)
        # Карточка с тонкой обводкой
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h,
        )
        card.adjustments[0] = 0.08
        card.fill.solid()
        card.fill.fore_color.rgb = NAVY_SOFT
        card.line.color.rgb = RGBColor(0x33, 0x45, 0x6B)
        card.line.width = Pt(1)
        card.shadow.inherit = False

        # Большая цифра
        add_text(
            slide, x, y + Inches(0.45), card_w, Inches(1.3),
            number, size=64, bold=True,
            color=ACCENT if i == 0 else WHITE,
            align=PP_ALIGN.CENTER,
        )
        # Лейбл
        add_text(
            slide, x, y + Inches(1.85), card_w, Inches(0.5),
            label, size=16, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER,
        )
        # Подпись
        add_text(
            slide, x + Inches(0.2), y + Inches(2.35),
            card_w - Inches(0.4), Inches(0.5),
            hint, size=11, color=RGBColor(0xCB, 0xD5, 0xE1),
            align=PP_ALIGN.CENTER, line_spacing=1.3,
        )

    add_text(
        slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.5),
        "В критериях курса 500 строк дают 15 баллов. У нас — почти "
        "в 9 раз больше.",
        size=14, color=RGBColor(0xCB, 0xD5, 0xE1),
        align=PP_ALIGN.CENTER,
    )

    add_footer(slide, "TripPlanner · защита проекта", dark=True)
    add_page_number(slide, 10, TOTAL_SLIDES, dark=True)


def slide_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG)
    _add_section_header(slide, "10", "Живое демо")

    # Левая колонка — URL и креды
    add_rounded_card(slide, Inches(0.5), Inches(1.7),
                     Inches(6.5), Inches(5.2))
    add_filled_rect(slide, Inches(0.5), Inches(1.7),
                    Inches(0.18), Inches(5.2), BLUE)

    add_text(
        slide, Inches(0.95), Inches(1.95), Inches(6), Inches(0.5),
        "ОТКРОЙТЕ В БРАУЗЕРЕ", size=12, bold=True, color=BLUE,
    )
    add_text(
        slide, Inches(0.95), Inches(2.4), Inches(6), Inches(0.8),
        "tripplanner-17ou.onrender.com",
        font="Consolas", size=22, bold=True, color=NAVY,
    )
    add_text(
        slide, Inches(0.95), Inches(3.3), Inches(6), Inches(0.4),
        "https://tripplanner-17ou.onrender.com",
        size=12, color=MUTED,
    )

    # Демо аккаунт
    add_filled_rect(slide, Inches(0.95), Inches(4.1),
                    Inches(6.0), Inches(1.4), NAVY)
    add_text(
        slide, Inches(1.15), Inches(4.2), Inches(5.7), Inches(0.4),
        "ДЕМО-АККАУНТ", size=11, bold=True,
        color=RGBColor(0xCB, 0xD5, 0xE1),
    )
    add_text(
        slide, Inches(1.15), Inches(4.6), Inches(5.7), Inches(0.4),
        "Логин: demo", font="Consolas", size=16, bold=True,
        color=WHITE,
    )
    add_text(
        slide, Inches(1.15), Inches(5.0), Inches(5.7), Inches(0.4),
        "Пароль: demo123", font="Consolas", size=16, bold=True,
        color=WHITE,
    )

    add_text(
        slide, Inches(0.95), Inches(5.7), Inches(6), Inches(0.4),
        "REST-проверка одним curl:",
        size=12, color=MUTED,
    )
    add_text(
        slide, Inches(0.95), Inches(6.1), Inches(6), Inches(0.7),
        "curl https://tripplanner-17ou.onrender.com/api/trips",
        font="Consolas", size=10, color=NAVY,
    )

    # Правая колонка — сценарий показа
    add_text(
        slide, Inches(7.4), Inches(1.85), Inches(5.5), Inches(0.5),
        "Что показать за 90 секунд",
        size=18, bold=True, color=NAVY,
    )

    script = [
        ("1", "Explore", "Открыть каталог публичных поездок."),
        ("2", "Открыть «Гастротур по Тбилиси»",
         "Карта с маркерами, расписание по дням."),
        ("3", "Лайк без перезагрузки",
         "Кнопка обновляет счётчик через AJAX."),
        ("4", "Войти как demo",
         "Создать новую точку через геокодер «Эрмитаж»."),
        ("5", "API-демо",
         "Открыть /api/trips → показать чистый JSON."),
    ]

    y0 = Inches(2.5)
    for index, (num, title, body) in enumerate(script):
        y = y0 + index * Inches(0.78)
        # Кружок с номером
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(7.4), y, Inches(0.42), Inches(0.42),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLUE
        circle.line.fill.background()
        circle.shadow.inherit = False
        add_text(
            slide, Inches(7.4), y, Inches(0.42), Inches(0.42),
            num, size=14, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide, Inches(7.95), y - Inches(0.02),
            Inches(5.0), Inches(0.35),
            title, size=14, bold=True, color=NAVY,
        )
        add_text(
            slide, Inches(7.95), y + Inches(0.32),
            Inches(5.0), Inches(0.4),
            body, size=11, color=MUTED,
        )

    add_footer(slide, "TripPlanner · защита проекта")
    add_page_number(slide, 11, TOTAL_SLIDES)


def slide_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, NAVY)

    # Боковая акцентная полоса
    add_filled_rect(slide, Inches(0), Inches(0), Inches(0.3),
                    Inches(7.5), ACCENT)

    add_text(
        slide, Inches(1.0), Inches(1.2), Inches(11), Inches(0.5),
        "СПАСИБО ЗА ВНИМАНИЕ",
        size=14, bold=True, color=ACCENT,
    )
    add_text(
        slide, Inches(1.0), Inches(1.8), Inches(11), Inches(1.4),
        "Готов ответить на ваши вопросы",
        size=44, bold=True, color=WHITE,
    )

    # Ссылки
    add_filled_rect(slide, Inches(1.0), Inches(3.6),
                    Inches(5.5), Inches(2.6), NAVY_SOFT)
    add_filled_rect(slide, Inches(1.0), Inches(3.6),
                    Inches(0.12), Inches(2.6), BLUE)
    add_text(
        slide, Inches(1.3), Inches(3.75), Inches(5), Inches(0.4),
        "САЙТ", size=11, bold=True, color=BLUE,
    )
    add_text(
        slide, Inches(1.3), Inches(4.15), Inches(5), Inches(0.5),
        "tripplanner-17ou.onrender.com",
        font="Consolas", size=15, bold=True, color=WHITE,
    )
    add_text(
        slide, Inches(1.3), Inches(4.8), Inches(5), Inches(0.4),
        "РЕПОЗИТОРИЙ", size=11, bold=True, color=BLUE,
    )
    add_text(
        slide, Inches(1.3), Inches(5.2), Inches(5), Inches(0.5),
        "github.com/rom0tod/travel",
        font="Consolas", size=15, bold=True, color=WHITE,
    )
    add_text(
        slide, Inches(1.3), Inches(5.85), Inches(5), Inches(0.4),
        "demo / demo123 — для входа на сайт",
        size=11, color=RGBColor(0xCB, 0xD5, 0xE1),
    )

    # Финальная цитата
    add_filled_rect(slide, Inches(7.0), Inches(3.6),
                    Inches(5.3), Inches(2.6), BLUE)
    add_text(
        slide, Inches(7.3), Inches(3.85), Inches(4.8), Inches(0.4),
        "В ОДНОМ ПРОЕКТЕ",
        size=11, bold=True, color=RGBColor(0xCB, 0xD5, 0xE1),
    )
    add_text(
        slide, Inches(7.3), Inches(4.25), Inches(4.8), Inches(2.0),
        "Flask · ORM · auth · WTForms · REST · сторонний API · "
        "файлы · Bootstrap · Leaflet · Postgres · gunicorn · "
        "Render.com",
        size=14, color=WHITE, line_spacing=1.35,
    )

    add_text(
        slide, Inches(1.0), Inches(6.75), Inches(12), Inches(0.4),
        "Роман Тодоров · Яндекс Лицей · 2026",
        size=12, color=RGBColor(0x94, 0xA3, 0xB8),
    )


def _add_section_header(slide, number, title, dark=False):
    """Маленькая шапка слайда: цифра + заголовок раздела."""
    text_color = WHITE if dark else NAVY
    add_text(
        slide, Inches(0.7), Inches(0.4), Inches(1.2), Inches(0.6),
        number, font=HEADER_FONT, size=42, bold=True, color=ACCENT,
    )
    add_text(
        slide, Inches(1.7), Inches(0.55), Inches(10), Inches(0.7),
        title, font=HEADER_FONT, size=32, bold=True,
        color=text_color,
    )
    # Тонкая линия-разделитель
    line_color = RGBColor(0x33, 0x45, 0x6B) if dark else BORDER
    add_filled_rect(slide, Inches(0.7), Inches(1.3),
                    Inches(11.9), Inches(0.02), line_color)


# --- Сборка ------------------------------------------------------------------

def build_presentation(output_path: Path) -> None:
    prs = Presentation()
    # 16:9 wide layout = 13.333 x 7.5 inches
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_problem(prs)
    slide_features(prs)
    slide_journey(prs)
    slide_architecture(prs)
    slide_data_model(prs)
    slide_tech_stack(prs)
    slide_api(prs)
    slide_highlights(prs)
    slide_numbers(prs)
    slide_demo(prs)
    slide_thanks(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Презентация сохранена: {output_path}")


if __name__ == "__main__":
    project_root = Path(__file__).resolve().parent.parent
    if str(project_root) not in sys.path:
        sys.path.insert(0, str(project_root))
    target = project_root / "docs" / "TripPlanner_presentation.pptx"
    build_presentation(target)
